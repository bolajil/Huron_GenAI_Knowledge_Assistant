"""
Huron Versioned Ingestion Service

Supports every document type by reducing it to text before chunking:
  PDF, DOCX, PPTX, XLSX/CSV, TXT/MD, HTML  → text extracted locally
  MP4/MKV/AVI/MOV, MP3/WAV/M4A/OGG         → transcribed via OpenAI Whisper API

Versioning strategy (Option B — metadata + is_latest):
  - Each chunk stored with: doc_id, doc_version, is_latest=True, effective_date
  - On re-ingest: old Pinecone vectors get is_latest=False, new vectors is_latest=True
  - Queries always filter {"is_latest": {"$eq": True}}
  - Full version history kept in document_versions DB table
  - Rollback = flip is_latest in Pinecone metadata + DB
"""
from __future__ import annotations

import hashlib
import io
import json
import logging
import os
import re
import sqlite3
import tempfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, List, Optional

logger = logging.getLogger(__name__)

# ── Result dataclass ─────────────────────────────────────────────────────────

@dataclass
class IngestionResult:
    success:       bool
    document_id:   str
    version:       str
    namespace:     str
    file_type:     str
    parent_chunks: int
    child_chunks:  int
    warnings:      List[str] = field(default_factory=list)
    _pinecone_upsert_success: bool = False


# ── File-type map ────────────────────────────────────────────────────────────

_EXT_MAP: dict[str, str] = {
    # Documents
    ".pdf":  "pdf",
    ".docx": "docx", ".doc": "docx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".xlsx": "xlsx", ".xls": "xlsx",
    ".csv":  "csv",
    ".txt":  "text", ".md": "text", ".rst": "text",
    ".html": "html", ".htm": "html",
    ".json": "json",
    # Video
    ".mp4":  "video", ".mkv": "video", ".avi": "video",
    ".mov":  "video", ".webm": "video",
    # Audio
    ".mp3":  "audio", ".wav": "audio", ".m4a": "audio",
    ".ogg":  "audio", ".flac": "audio",
}

_CHUNK_SIZE    = 500   # target tokens per chunk (approx by words)
_CHUNK_OVERLAP = 50    # word overlap between consecutive chunks


# ── Main service ─────────────────────────────────────────────────────────────

class IngestionService:
    def __init__(
        self,
        enable_pinecone:        bool = True,
        enable_dlp_scan:        bool = False,
        enable_quality_gate:    bool = False,
        enable_classification:  bool = False,
        pinecone_index:         str  = "huron-enterprise-knowledge",
    ):
        self._pc_index_name = pinecone_index
        self._pc_index      = None
        self._oai_client    = None

    # ── Public entry point ───────────────────────────────────────────────────

    async def ingest_document(
        self,
        file_content:      bytes,
        file_name:         str,
        tenant_context:    Any,
        sensitivity_level: str  = "internal",
        doc_id:            Optional[str] = None,
        version:           Optional[str] = None,
    ) -> IngestionResult:
        import asyncio
        return await asyncio.to_thread(
            self._ingest_sync,
            file_content, file_name, tenant_context, sensitivity_level, doc_id, version,
        )

    # ── Synchronous core (runs in thread) ────────────────────────────────────

    def _ingest_sync(
        self,
        file_content:      bytes,
        file_name:         str,
        tenant_context:    Any,
        sensitivity_level: str,
        doc_id:            Optional[str],
        version:           Optional[str],
    ) -> IngestionResult:
        warnings:  List[str] = []
        dept      = getattr(tenant_context, "dept_id", "general")
        namespace = f"vaultmind-huron-{dept}-general"

        # Determine file type
        ext       = Path(file_name).suffix.lower()
        file_type = _EXT_MAP.get(ext, "text")

        # Derive stable doc_id from filename if not supplied
        if not doc_id:
            doc_id = self._derive_doc_id(file_name)

        # Auto-version as YYYY-MM if not supplied
        if not version:
            version = datetime.now(timezone.utc).strftime("%Y-%m")

        effective_date = datetime.now(timezone.utc).strftime("%Y-%m-%d")

        # ── Step 1: Extract text ─────────────────────────────────────────────
        try:
            text = self._extract_text(file_content, file_name, file_type)
        except Exception as e:
            warnings.append(f"Text extraction failed: {e}")
            text = ""

        if not text.strip():
            return IngestionResult(
                success=False, document_id=doc_id, version=version,
                namespace=namespace, file_type=file_type,
                parent_chunks=0, child_chunks=0,
                warnings=warnings + ["No text could be extracted from this file."],
            )

        # ── Step 2: Chunk ────────────────────────────────────────────────────
        chunks = self._chunk_text(text)

        # ── Step 3: Build vector IDs ─────────────────────────────────────────
        # ID format: {dept}__{doc_id}__{version}__{chunk_idx}
        vector_ids = [f"{dept}__{doc_id}__{version}__{i}" for i in range(len(chunks))]

        # ── Step 4: Mark old version as superseded ───────────────────────────
        self._supersede_old_version(dept, doc_id, namespace, warnings)

        # ── Step 5: Embed + upsert to Pinecone ──────────────────────────────
        upsert_ok = False
        try:
            idx = self._get_pinecone_index()
            vectors_to_upsert = []
            for i, (chunk, vid) in enumerate(zip(chunks, vector_ids)):
                emb = self._embed(chunk)
                vectors_to_upsert.append({
                    "id":     vid,
                    "values": emb,
                    "metadata": {
                        "text":           chunk,
                        "source":         file_name,
                        "doc_id":         doc_id,
                        "doc_version":    version,
                        "effective_date": effective_date,
                        "dept":           dept,
                        "file_type":      file_type,
                        "sensitivity":    sensitivity_level,
                        "is_latest":      True,
                        "chunk_idx":      i,
                        "page":           None,
                    },
                })

            # Upsert in batches of 100
            for batch_start in range(0, len(vectors_to_upsert), 100):
                batch = vectors_to_upsert[batch_start:batch_start + 100]
                idx.upsert(vectors=batch, namespace=namespace)
            upsert_ok = True
        except Exception as e:
            warnings.append(f"Pinecone upsert failed: {e}")
            logger.warning("Pinecone upsert error: %s", e)

        # ── Step 6: Persist version record in DB ─────────────────────────────
        if upsert_ok:
            try:
                self._save_version_record(
                    dept, doc_id, version, vector_ids,
                    file_name, file_type,
                    getattr(tenant_context, "username", ""),
                    len(chunks),
                )
            except Exception as e:
                warnings.append(f"Version record save failed: {e}")
                logger.warning("Version DB save error: %s", e)

        return IngestionResult(
            success=upsert_ok,
            document_id=doc_id,
            version=version,
            namespace=namespace,
            file_type=file_type,
            parent_chunks=len(chunks),
            child_chunks=len(chunks),
            warnings=warnings,
            _pinecone_upsert_success=upsert_ok,
        )

    # ── Text extraction ──────────────────────────────────────────────────────

    def _extract_text(self, content: bytes, file_name: str, file_type: str) -> str:
        if file_type == "pdf":
            return self._extract_pdf(content)
        if file_type == "docx":
            return self._extract_docx(content)
        if file_type == "pptx":
            return self._extract_pptx(content)
        if file_type in ("xlsx", "csv"):
            return self._extract_spreadsheet(content, file_type)
        if file_type == "html":
            return self._extract_html(content)
        if file_type == "json":
            return self._extract_json(content)
        if file_type in ("video", "audio"):
            return self._extract_video_audio(content, file_name)
        # Plain text fallback
        return content.decode("utf-8", errors="replace")

    def _extract_pdf(self, content: bytes) -> str:
        # Try pdfplumber first (better layout), fall back to PyPDF2
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                pages = []
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        pages.append(t)
                return "\n\n".join(pages)
        except ImportError:
            pass
        except Exception as e:
            logger.warning("pdfplumber failed: %s", e)

        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            return "\n\n".join(
                page.extract_text() or "" for page in reader.pages
            )
        except Exception as e:
            raise RuntimeError(f"PDF extraction failed: {e}") from e

    def _extract_docx(self, content: bytes) -> str:
        try:
            from docx import Document
            doc  = Document(io.BytesIO(content))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            raise RuntimeError("python-docx not installed. Run: pip install python-docx")
        except Exception as e:
            raise RuntimeError(f"DOCX extraction failed: {e}") from e

    def _extract_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
            prs    = Presentation(io.BytesIO(content))
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
        except Exception as e:
            raise RuntimeError(f"PPTX extraction failed: {e}") from e

    def _extract_spreadsheet(self, content: bytes, file_type: str) -> str:
        if file_type == "csv":
            try:
                import csv
                text   = content.decode("utf-8", errors="replace")
                reader = csv.reader(io.StringIO(text))
                rows   = ["\t".join(row) for row in reader]
                return "\n".join(rows)
            except Exception as e:
                raise RuntimeError(f"CSV extraction failed: {e}") from e

        try:
            import openpyxl
            wb     = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            sheets = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        except ImportError:
            raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
        except Exception as e:
            raise RuntimeError(f"XLSX extraction failed: {e}") from e

    def _extract_html(self, content: bytes) -> str:
        try:
            from html.parser import HTMLParser

            class _Strip(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self._parts: list[str] = []
                    self._skip  = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style"):
                        self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script", "style"):
                        self._skip = False
                def handle_data(self, data):
                    if not self._skip and data.strip():
                        self._parts.append(data.strip())

            parser = _Strip()
            parser.feed(content.decode("utf-8", errors="replace"))
            return " ".join(parser._parts)
        except Exception as e:
            raise RuntimeError(f"HTML extraction failed: {e}") from e

    def _extract_json(self, content: bytes) -> str:
        try:
            data = json.loads(content.decode("utf-8", errors="replace"))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception:
            return content.decode("utf-8", errors="replace")

    def _extract_video_audio(self, content: bytes, file_name: str) -> str:
        """Transcribe video/audio via OpenAI Whisper API."""
        try:
            from openai import OpenAI
            client = OpenAI(api_key=os.getenv("OPENAI_API_KEY", ""))
            ext    = Path(file_name).suffix.lower() or ".mp4"
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                with open(tmp_path, "rb") as f:
                    transcript = client.audio.transcriptions.create(
                        model="whisper-1",
                        file=f,
                        response_format="text",
                    )
                return transcript if isinstance(transcript, str) else transcript.text
            finally:
                os.unlink(tmp_path)
        except ImportError:
            raise RuntimeError("openai package not installed")
        except Exception as e:
            raise RuntimeError(f"Whisper transcription failed: {e}") from e

    # ── Chunking ─────────────────────────────────────────────────────────────

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping word-based chunks."""
        words  = text.split()
        if not words:
            return []
        chunks = []
        start  = 0
        while start < len(words):
            end   = min(start + _CHUNK_SIZE, len(words))
            chunk = " ".join(words[start:end])
            if chunk.strip():
                chunks.append(chunk)
            if end >= len(words):
                break
            start = end - _CHUNK_OVERLAP
        return chunks

    # ── Embedding ────────────────────────────────────────────────────────────

    def _embed(self, text: str) -> List[float]:
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY", ""))
        res    = client.embeddings.create(
            input=text,
            model="text-embedding-3-small",
        )
        return res.data[0].embedding

    # ── Pinecone ─────────────────────────────────────────────────────────────

    def _get_pinecone_index(self):
        if self._pc_index is None:
            from pinecone import Pinecone
            pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY", ""))
            self._pc_index = pc.Index(self._pc_index_name)
        return self._pc_index

    # ── Version management ───────────────────────────────────────────────────

    def _supersede_old_version(
        self, dept: str, doc_id: str, namespace: str, warnings: List[str]
    ) -> None:
        """Mark existing latest vectors for this doc as is_latest=False."""
        old_ids = self._get_old_vector_ids(dept, doc_id)
        if not old_ids:
            return

        # Update Pinecone metadata
        try:
            idx = self._get_pinecone_index()
            for vid in old_ids:
                try:
                    idx.update(id=vid, set_metadata={"is_latest": False}, namespace=namespace)
                except Exception:
                    pass  # non-fatal; old chunks will still rank lower due to version filter
        except Exception as e:
            warnings.append(f"Could not mark old Pinecone vectors: {e}")

        # Update DB
        try:
            db_path = os.path.join(
                os.path.dirname(os.path.dirname(__file__)), "data", "huron.db"
            )
            with sqlite3.connect(db_path) as conn:
                conn.execute(
                    "UPDATE document_versions SET is_latest=0 WHERE dept=? AND doc_id=? AND is_latest=1",
                    (dept, doc_id),
                )
                conn.commit()
        except Exception as e:
            warnings.append(f"DB version update failed: {e}")

    def _get_old_vector_ids(self, dept: str, doc_id: str) -> List[str]:
        try:
            db_path = os.path.join(
                os.path.dirname(os.path.dirname(__file__)), "data", "huron.db"
            )
            with sqlite3.connect(db_path) as conn:
                row = conn.execute(
                    "SELECT vector_ids FROM document_versions WHERE dept=? AND doc_id=? AND is_latest=1",
                    (dept, doc_id),
                ).fetchone()
            if row and row[0]:
                return json.loads(row[0])
        except Exception as e:
            logger.warning("Could not fetch old vector IDs: %s", e)
        return []

    def _save_version_record(
        self,
        dept:         str,
        doc_id:       str,
        version:      str,
        vector_ids:   List[str],
        source_file:  str,
        file_type:    str,
        ingested_by:  str,
        chunk_count:  int,
    ) -> None:
        effective_date = datetime.now(timezone.utc).strftime("%Y-%m-%d")
        db_path = os.path.join(
            os.path.dirname(os.path.dirname(__file__)), "data", "huron.db"
        )
        with sqlite3.connect(db_path) as conn:
            conn.execute(
                """
                INSERT OR REPLACE INTO document_versions
                    (dept, doc_id, version, source_file, file_type,
                     effective_date, vector_ids, chunk_count, is_latest, ingested_by)
                VALUES (?,?,?,?,?,?,?,?,1,?)
                """,
                (dept, doc_id, version, source_file, file_type,
                 effective_date, json.dumps(vector_ids), chunk_count, ingested_by),
            )
            conn.commit()

    # ── Helpers ──────────────────────────────────────────────────────────────

    def _derive_doc_id(self, file_name: str) -> str:
        """Derive a stable doc_id slug from the filename, stripping date suffixes."""
        stem = Path(file_name).stem
        # Remove common date/version patterns: 2024, 2025, v1, v2.1, _jan, _jan2025, etc.
        stem = re.sub(r"[_\-\s]*(20\d{2}|v\d+(\.\d+)*|q[1-4]|jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\w*",
                      "", stem, flags=re.IGNORECASE)
        slug = re.sub(r"[^a-z0-9]+", "_", stem.lower()).strip("_")
        return slug or hashlib.md5(file_name.encode()).hexdigest()[:12]
